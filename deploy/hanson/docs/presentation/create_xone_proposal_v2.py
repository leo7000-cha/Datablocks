#!/usr/bin/env python3
"""
X-One 통합 플랫폼 개인정보 거버넌스 제안서 v2
JB우리캐피탈 프로토타입 설명회용
- 수직 중앙 정렬 적용
- 번호 줄바꿈 수정
- 폰트/비주얼 개선
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Color Palette ──
C_NAVY      = RGBColor(0x1B, 0x2A, 0x4A)
C_DEEP_BLUE = RGBColor(0x1E, 0x3A, 0x5F)
C_TEAL      = RGBColor(0x00, 0x7B, 0x8A)
C_ACCENT    = RGBColor(0xE8, 0x6C, 0x00)
C_ACCENT2   = RGBColor(0x2E, 0x86, 0xAB)
C_GREEN     = RGBColor(0x00, 0x8A, 0x5E)
C_RED       = RGBColor(0xC0, 0x39, 0x2B)
C_LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK     = RGBColor(0x1A, 0x1A, 0x2E)
C_GRAY      = RGBColor(0x6B, 0x72, 0x80)
C_LIGHT_GRAY= RGBColor(0xE5, 0xE7, 0xEB)
C_VERY_LIGHT= RGBColor(0xF8, 0xFA, 0xFC)
C_PURPLE    = RGBColor(0x6C, 0x3E, 0xB2)
C_GOLD      = RGBColor(0xD4, 0x9B, 0x10)

prs = Presentation()
prs.slide_width = Emu(12192000)   # 16:9
prs.slide_height = Emu(6858000)

SLIDE_W = 12192000
SLIDE_H = 6858000

# ── Helper Functions ──

def _set_vcenter(text_frame):
    """텍스트 프레임 수직 중앙 정렬"""
    txBody = text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')

def _set_vtop(text_frame):
    """텍스트 프레임 수직 상단 정렬"""
    txBody = text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 't')

def _set_no_autofit(text_frame):
    """자동 줄임 끄기"""
    txBody = text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('wrap', 'square')
        # Remove spAutoFit if exists
        for child in bodyPr:
            if child.tag.endswith('spAutoFit') or child.tag.endswith('normAutofit'):
                bodyPr.remove(child)

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=14, bold=False,
             color=C_BLACK, align=PP_ALIGN.LEFT, valign='middle', font_name='맑은 고딕',
             wrap=True, line_spacing=1.15):
    """텍스트 박스 추가 — 수직 정렬 기본 middle"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    if line_spacing:
        p.line_spacing = line_spacing

    if valign == 'middle':
        _set_vcenter(tf)
    elif valign == 'top':
        _set_vtop(tf)
    _set_no_autofit(tf)
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=11,
                  color=C_BLACK, bold=False, align=PP_ALIGN.LEFT, valign='top',
                  font_name='맑은 고딕', bullet='·', spacing=1.3):
    """여러 줄 텍스트 (bullet 포함)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f'{bullet}  {line}' if bullet else line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        p.line_spacing = spacing

    if valign == 'middle':
        _set_vcenter(tf)
    elif valign == 'top':
        _set_vtop(tf)
    _set_no_autofit(tf)
    return txBox

def add_shape_text(slide, left, top, width, height, text, font_size=14, bold=False,
                   color=C_WHITE, fill_color=None, line_color=None, line_width=None,
                   align=PP_ALIGN.CENTER, font_name='맑은 고딕', shape_type='rect'):
    """텍스트가 포함된 도형 — 수직 중앙 정렬 자동 적용"""
    if shape_type == 'rounded':
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    elif shape_type == 'oval':
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    _set_vcenter(tf)
    return shape

def add_page_number(slide, num):
    add_text(slide, Cm(31), Cm(17.5), Cm(2.5), Cm(0.8),
             f'{num:02d}', font_size=9, color=C_GRAY, align=PP_ALIGN.RIGHT)

def add_bottom_bar(slide, color=C_NAVY):
    add_shape(slide, 0, Emu(SLIDE_H - Cm(0.35).emu), Emu(SLIDE_W), Cm(0.35), fill_color=color)

def add_top_accent(slide, color=C_TEAL):
    add_shape(slide, 0, 0, Emu(SLIDE_W), Cm(0.12), fill_color=color)

def add_section_header(slide, section_num, title, subtitle="", header_color=C_NAVY):
    """섹션 헤더 바 — 번호 줄바꿈 방지"""
    add_shape(slide, 0, 0, Emu(SLIDE_W), Cm(2.3), fill_color=header_color)
    # 번호: 충분한 너비로
    add_text(slide, Cm(1.2), Cm(0.15), Cm(3), Cm(2),
             f'{section_num:02d}', font_size=28, bold=True, color=C_ACCENT, wrap=False)
    # 제목
    add_text(slide, Cm(4), Cm(0.2), Cm(28), Cm(1.3),
             title, font_size=20, bold=True, color=C_WHITE, valign='middle')
    if subtitle:
        add_text(slide, Cm(4), Cm(1.45), Cm(28), Cm(0.7),
                 subtitle, font_size=11, color=RGBColor(0xB0, 0xC4, 0xD8), valign='middle')

def add_key_point(slide, text, top=Cm(17), color=C_DEEP_BLUE):
    add_shape_text(slide, Cm(1), top, Cm(31.5), Cm(0.9),
                   f'  {text}', font_size=11, bold=True, color=C_WHITE,
                   fill_color=color, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════
# ═══ Slide 1: 표지 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, Emu(SLIDE_W), Emu(SLIDE_H), fill_color=C_NAVY)
add_shape(slide, 0, 0, Emu(SLIDE_W), Cm(0.2), fill_color=C_TEAL)
add_shape(slide, Cm(2.5), Cm(3), Cm(0.15), Cm(11), fill_color=C_TEAL)

add_text(slide, Cm(3.5), Cm(3.5), Cm(28), Cm(1),
         'ENTERPRISE PRIVACY GOVERNANCE PLATFORM', font_size=13,
         color=C_TEAL, valign='middle')
add_text(slide, Cm(3.5), Cm(5), Cm(20), Cm(2.5),
         'X-One', font_size=60, bold=True, color=C_WHITE, valign='middle')
add_text(slide, Cm(3.5), Cm(8), Cm(20), Cm(1.2),
         'All Data.  One Platform.', font_size=22, color=RGBColor(0x8A, 0xB4, 0xD0))
add_text(slide, Cm(3.5), Cm(9.8), Cm(28), Cm(1.2),
         '전사 개인정보 거버넌스 통합 플랫폼', font_size=18, bold=True, color=C_WHITE)
add_text(slide, Cm(3.5), Cm(11.2), Cm(28), Cm(1),
         '"잘 관리하고 있습니다"를 시스템으로 증명할 수 있는 유일한 체계', font_size=14, color=RGBColor(0xC0, 0xD0, 0xE0))

add_shape(slide, Cm(3.5), Cm(14.5), Cm(12), Cm(0.05), fill_color=C_TEAL)
add_text(slide, Cm(3.5), Cm(14.8), Cm(20), Cm(0.9),
         'JB우리캐피탈  솔루션 제안', font_size=16, bold=True, color=C_WHITE)
add_text(slide, Cm(3.5), Cm(15.9), Cm(20), Cm(0.7),
         '2026년 4월  |  데이터블록스 (Datablocks)', font_size=11, color=C_GRAY)
add_page_number(slide, 1)


# ═══════════════════════════════════════════════════════
# ═══ Slide 2: 목차 (TABLE OF CONTENTS) ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, Emu(SLIDE_W), Emu(SLIDE_H), fill_color=C_WHITE)
add_top_accent(slide)
add_shape(slide, 0, 0, Emu(SLIDE_W), Cm(2.3), fill_color=C_NAVY)
add_text(slide, Cm(1.5), Cm(0.3), Cm(28), Cm(1.8),
         'CONTENTS', font_size=24, bold=True, color=C_WHITE, valign='middle')
add_text(slide, Cm(20), Cm(0.5), Cm(12), Cm(1.5),
         '목차', font_size=13, color=RGBColor(0xB0, 0xC4, 0xD8), align=PP_ALIGN.RIGHT, valign='middle')

# 좌측 수직 라인
add_shape(slide, Cm(2.5), Cm(3.2), Cm(0.1), Cm(14.5), fill_color=C_TEAL)

toc_sections = [
    ("01", "도입 배경",
     [("핵심 질문 — 개인정보 관리 현실", "03"),
      ("규제 환경 — 2026 개인정보보호법 개정", "04"),
      ("리스크 시뮬레이션 / 기존 방식 한계", "05")],
     C_NAVY),
    ("02", "X-One 플랫폼",
     [("All-in-One 거버넌스 접근 방식", "07"),
      ("플랫폼 아키텍처 / X-Scan", "09"),
      ("X-Purge 파기 자동화", "11")],
     C_TEAL),
    ("03", "X-Audit 접속기록 · 소명",
     [("5W1H 접속기록 수집 / 로그 전략", "14"),
      ("이상행위 탐지 7가지 규칙", "16"),
      ("소명 프로세스 / 위변조 방지", "17"),
      ("경쟁 솔루션 대비 차별화", "19")],
     C_ACCENT),
    ("04", "도입 효과 & 로드맵",
     [("도입 효과 & ROI", "23"),
      ("JB우리캐피탈 도입 로드맵 (8주)", "25"),
      ("거버넌스 완성 로드맵", "26")],
     C_DEEP_BLUE),
    ("05", "AI 로드맵 & 레퍼런스",
     [("AI 로드맵 — X-One Intelligence", "27"),
      ("도입 레퍼런스 (금융권)", "28")],
     C_PURPLE),
]

# 목차 레이아웃: 5개 섹션을 균등 배치
section_height = Cm(2.6)
item_row_h = Cm(0.6)

for i, (num, title, items, color) in enumerate(toc_sections):
    y = Emu(Cm(3.2).emu + int(i * section_height.emu))

    # 섹션 번호 (원형)
    add_shape_text(slide, Cm(3.5), Emu(y.emu + Cm(0.15).emu), Cm(1.6), Cm(1.6),
                   num, font_size=13, bold=True, color=C_WHITE, fill_color=color, shape_type='oval')

    # 섹션 제목
    add_text(slide, Cm(5.8), y, Cm(9), Cm(1.8),
             title, font_size=15, bold=True, color=color, valign='middle')

    # 하위 항목 — 오른쪽 영역
    for j, (item_text, page) in enumerate(items):
        item_y = Emu(y.emu + int(j * item_row_h.emu))
        # 항목명
        add_text(slide, Cm(16), item_y, Cm(14), item_row_h,
                 item_text, font_size=10, color=C_GRAY, valign='middle')
        # 페이지 번호
        add_text(slide, Cm(31), item_y, Cm(1.8), item_row_h,
                 page, font_size=10, bold=True, color=color, align=PP_ALIGN.RIGHT, valign='middle')

    # 구분선 (마지막 제외)
    if i < len(toc_sections) - 1:
        line_y = Emu(y.emu + int(max(len(items), 2) * item_row_h.emu) + Cm(0.3).emu)
        add_shape(slide, Cm(3.5), line_y, Cm(29), Cm(0.03), fill_color=C_LIGHT_GRAY)

add_bottom_bar(slide)
add_page_number(slide, 2)


# ═══════════════════════════════════════════════════════
# ═══ Slide 3: 핵심 질문 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 1, 'THE CORE QUESTION',
                   '"잘 관리하고 있습니다" — 정말 그렇게 말할 수 있습니까?')

warnings = [
    ("개인정보가 어디에 있는지 모릅니다",
     ["전사 DB에 산재된 개인정보 현황을 정확히 인지하고 있습니까?",
      "사각지대 없이 모든 테이블/컬럼을 파악하고 있습니까?"]),
    ("링크 단절 ≠ 파기",
     ["고객 마스터 테이블의 링크만 끊는 방식은 '파기'로 인정되지 않습니다.",
      "데이터는 여전히 DB에 존재하며, 유출 시 '파기 미이행'으로 처벌됩니다."]),
    ("누가 언제 왜 접근했는지 모릅니다",
     ["접속기록을 저장만 하고, 이상행위를 탐지하고 소명을 받는 체계가",
      "없다면 '안전조치 의무 위반'입니다."]),
    ("테스트 환경에 실제 개인정보가 있습니다",
     ["개발/테스트 DB에 운영 데이터를 그대로 복사해서 사용하고",
      "있지 않습니까? 이것도 개인정보 보호법 위반입니다."]),
]

for i, (title, items) in enumerate(warnings):
    x = Cm(1.5) if i < 2 else Cm(17)
    y = Cm(3.2) if i % 2 == 0 else Cm(10.2)

    card = add_rounded_rect(slide, x, y, Cm(15), Cm(6), fill_color=C_VERY_LIGHT, line_color=C_LIGHT_GRAY, line_width=Pt(1))
    add_shape(slide, x, y, Cm(15), Cm(0.15), fill_color=C_RED)

    add_text(slide, Emu(x.emu + Cm(1).emu), Emu(y.emu + Cm(0.7).emu), Cm(13.5), Cm(1.2),
             title, font_size=16, bold=True, color=C_RED, valign='middle')
    add_multiline(slide, Emu(x.emu + Cm(1).emu), Emu(y.emu + Cm(2.5).emu), Cm(13.5), Cm(3),
                  items, font_size=12, color=C_BLACK, bullet='')

add_key_point(slide, '전사 개인정보의 인지 → 파기 → 감시 → 안전한 활용까지, 통합 거버넌스 체계가 필요합니다.')
add_bottom_bar(slide)
add_page_number(slide, 3)


# ═══════════════════════════════════════════════════════
# ═══ Slide 3: 규제 환경 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 2, '규제 환경 급변 — 2026년 개인정보보호법 개정',
                   '2026.2.12 국회 본회의 통과  |  관리 체계 부재 자체가 처벌 대상')

cards_data = [
    ("매출액 10%", "과징금 상한", ["관련매출 3% → 전체매출 3%", "중대·반복 시 10%", "모든 개인정보 처리자 확대 적용"], C_RED),
    ("3천만원", "항목별 과태료", ["파기 미이행 — 3천만원", "접속기록 미관리 — 3천만원", "안전조치 미흡 — 3천만원", "분리보관 위반 — 1천만원"], C_ACCENT),
    ("CPO·대표자", "직접 책임 명시", ["대표자 최종 책임 명시", "CPO 이사회 보고 의무", "위탁사 감독 소홀 시 1차 책임"], C_DEEP_BLUE),
]

for i, (big_num, label, items, color) in enumerate(cards_data):
    x = Cm(1.5 + i * 10.8)
    add_rounded_rect(slide, x, Cm(3.2), Cm(10), Cm(10.5), fill_color=C_WHITE, line_color=C_LIGHT_GRAY, line_width=Pt(1))
    add_shape(slide, x, Cm(3.2), Cm(10), Cm(0.15), fill_color=color)

    add_text(slide, Emu(x.emu + Cm(0.5).emu), Cm(3.8), Cm(9), Cm(2),
             big_num, font_size=30, bold=True, color=color, align=PP_ALIGN.CENTER, valign='middle')
    add_text(slide, Emu(x.emu + Cm(0.5).emu), Cm(5.8), Cm(9), Cm(1),
             label, font_size=14, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER, valign='middle')

    add_multiline(slide, Emu(x.emu + Cm(1.2).emu), Cm(7.5), Cm(8), Cm(5.5),
                  items, font_size=12, color=C_BLACK)

# 법조항
add_text(slide, Cm(1.5), Cm(14.2), Cm(31), Cm(0.7),
         '개인정보보호법 제21조: 보유기간 경과 시 지체 없이 파기  |  제29조: 안전조치 의무',
         font_size=10, color=C_GRAY)
add_text(slide, Cm(1.5), Cm(15), Cm(31), Cm(0.7),
         '안전성 확보조치 제8조: 접속기록 최소 1년(금융사 5년)  |  월 1회 점검  |  위변조 방지',
         font_size=10, color=C_GRAY)

add_key_point(slide, '핵심: 솔루션 도입·투자 증빙은 과징금 감경 실질 수단 (개정법 명시)')
add_bottom_bar(slide)
add_page_number(slide, 4)


# ═══════════════════════════════════════════════════════
# ═══ Slide 4: 리스크 시뮬레이션 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 3, '리스크 시뮬레이션',
                   '시나리오: 캐피탈사 — 개인정보 관리 체계 부재 적발 시')

steps = [
    ("STEP 1", "파기 미이행 적발", "개보위 직권조사\n또는 민원 접수", "과태료\n3천만원", C_ACCENT),
    ("STEP 2", "접속기록 관리 부재", "접속기록 미보관\n점검·소명 미이행", "추가 과태료\n3천만원+", C_RED),
    ("STEP 3", "유출 사고 연계", "안전조치 소홀\n+ 유출 발생", "전체 매출\n3%", RGBColor(0xA0, 0x20, 0x20)),
    ("STEP 4", "반복·중대 위반", "3년 내 반복\n1천만명 피해", "전체 매출\n10%", RGBColor(0x80, 0x00, 0x00)),
]

for i, (step, title, desc, amount, color) in enumerate(steps):
    x = Cm(1.2 + i * 8.2)
    card = add_rounded_rect(slide, x, Cm(3.5), Cm(7.5), Cm(10), fill_color=C_WHITE, line_color=color, line_width=Pt(2))

    add_shape_text(slide, x, Cm(3.5), Cm(7.5), Cm(1.5),
                   f'{step}  {title}', font_size=12, bold=True, color=C_WHITE, fill_color=color)

    add_text(slide, Emu(x.emu + Cm(0.5).emu), Cm(5.8), Cm(6.5), Cm(2.5),
             desc, font_size=12, color=C_BLACK, align=PP_ALIGN.CENTER, valign='middle')

    add_shape(slide, Emu(x.emu + Cm(1.2).emu), Cm(9), Cm(5.1), Cm(0.05), fill_color=C_LIGHT_GRAY)
    add_text(slide, Emu(x.emu + Cm(0.5).emu), Cm(9.5), Cm(6.5), Cm(3),
             amount, font_size=20, bold=True, color=color, align=PP_ALIGN.CENTER, valign='middle')

    if i < 3:
        add_text(slide, Cm(8.8 + i * 8.2), Cm(7.5), Cm(1.5), Cm(1.5),
                 '→', font_size=26, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER, valign='middle')

warnings_text = [
    '파기 미이행 + 접속기록 미보관 = 과태료 중복 부과 (누적 6천만원+)',
    '대표자·CPO 직접 책임 — 임원진 법적 리스크로 격상',
    '솔루션 도입·투자 증빙 = 과징금 감경 실질 수단',
]
for j, w in enumerate(warnings_text):
    add_text(slide, Cm(1.5), Cm(14.2 + j * 0.85), Cm(31), Cm(0.8),
             f'  ▸  {w}', font_size=11, bold=True, color=C_RED)

add_bottom_bar(slide)
add_page_number(slide, 5)


# ═══════════════════════════════════════════════════════
# ═══ Slide 5: 기존 방식의 한계 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 4, '기존 방식의 구조적 한계',
                   '개별 대응으로는 증명이 불가능합니다')

limits = [
    ("개인정보 인지", ["PII 현황 파악 불가 — 전사 DB 산재", "수작업 메타 관리 — Excel, 최신화 불가"], C_ACCENT2),
    ("파기", ["SQL 직접 작성 의존 — DBA 퇴사 시 유실", "연관추적 실패 — 불완전 파기, 잔존 리스크"], C_RED),
    ("접근이력", ["접속기록 미관리 — 수집/저장 체계 없음", "소명체계 부재 — 이상행위 시 절차 없음"], C_ACCENT),
    ("테스트데이터", ["운영데이터 직접 복사 — 개발DB에 실 PII", "마스킹 수작업 — 컬럼별 개별 SQL 변환"], C_TEAL),
]

for i, (title, items, color) in enumerate(limits):
    x = Cm(1.5 + i * 8)
    add_rounded_rect(slide, x, Cm(3.5), Cm(7.2), Cm(10.5), fill_color=C_VERY_LIGHT, line_color=color, line_width=Pt(1.5))

    add_shape_text(slide, x, Cm(3.5), Cm(7.2), Cm(1.6),
                   title, font_size=15, bold=True, color=C_WHITE, fill_color=color)

    for j, item in enumerate(items):
        add_text(slide, Emu(x.emu + Cm(0.5).emu), Cm(6 + j * 3.2), Cm(0.8), Cm(0.8),
                 '✕', font_size=18, bold=True, color=C_RED, align=PP_ALIGN.CENTER, valign='middle')
        add_text(slide, Emu(x.emu + Cm(1.5).emu), Cm(6 + j * 3.2), Cm(5.2), Cm(2.5),
                 item, font_size=11, color=C_BLACK, valign='middle')

add_key_point(slide, '필요한 것은 개별 솔루션이 아니라, 하나의 플랫폼에서 통합 관리하는 거버넌스 체계입니다.')
add_bottom_bar(slide)
add_page_number(slide, 6)


# ═══════════════════════════════════════════════════════
# ═══ Slide 6: X-One 접근 방식 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 5, 'X-One의 접근 방식: All-in-One 거버넌스',
                   '개인정보 관리의 4대 영역을 하나의 플랫폼에서 통합합니다')

modules = [
    ("01", "인지  DISCOVER", "X-Scan", "어디에 어떤 PII가\n있는가", C_ACCENT2),
    ("02", "통제  GOVERN", "X-Purge", "보유기간 경과 시\n완전 파기", C_RED),
    ("03", "감시  MONITOR", "X-Audit", "누가 언제 왜\n접근했는가", C_ACCENT),
    ("04", "활용  UTILIZE", "X-Gen", "안전한 테스트\n데이터 생성", C_TEAL),
]

for i, (num, label, name, desc, color) in enumerate(modules):
    x = Cm(1.5 + i * 8)
    add_rounded_rect(slide, x, Cm(3.5), Cm(7.2), Cm(8.5), fill_color=C_WHITE, line_color=color, line_width=Pt(2))

    add_shape_text(slide, Emu(x.emu + Cm(2.6).emu), Cm(4), Cm(2), Cm(2),
                   num, font_size=18, bold=True, color=C_WHITE, fill_color=color, shape_type='oval')

    add_text(slide, x, Cm(6.5), Cm(7.2), Cm(0.8),
             label, font_size=11, bold=True, color=color, align=PP_ALIGN.CENTER, valign='middle')
    add_text(slide, x, Cm(7.5), Cm(7.2), Cm(1.3),
             name, font_size=24, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER, valign='middle')
    add_text(slide, Emu(x.emu + Cm(0.5).emu), Cm(9.3), Cm(6.2), Cm(2),
             desc, font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER, valign='middle')

# Hub bar
add_shape_text(slide, Cm(1.5), Cm(12.8), Cm(31), Cm(1),
               'X-One Hub', font_size=18, bold=True, color=C_TEAL, fill_color=C_NAVY)
hub_items = ['PII 메타데이터 전 모듈 공유', '통합 CPO 대시보드', '단일 인프라 · 단일 관리 콘솔', '감사 리포트 통합 생성']
for i, item in enumerate(hub_items):
    add_text(slide, Cm(2 + i * 7.8), Cm(14), Cm(7.5), Cm(0.7),
             f'·  {item}', font_size=10, color=C_WHITE)

add_key_point(slide, 'X-One은 4단계 체계를 하나의 플랫폼으로 구현하는 All-in-One 개인정보 거버넌스 솔루션입니다.', top=Cm(16.5))
add_bottom_bar(slide)
add_page_number(slide, 7)


# ═══════════════════════════════════════════════════════
# ═══ Slide 7: All-in-One 비교 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 6, '왜 All-in-One인가 — 개별 솔루션 도입 대비 압도적 우위')

# Header
add_shape_text(slide, Cm(1.5), Cm(3.2), Cm(7), Cm(1.2),
               '구분', font_size=13, bold=True, color=C_WHITE, fill_color=C_GRAY)
add_shape_text(slide, Cm(8.5), Cm(3.2), Cm(12), Cm(1.2),
               'X-One (All-in-One)', font_size=13, bold=True, color=C_WHITE, fill_color=C_TEAL)
add_shape_text(slide, Cm(20.5), Cm(3.2), Cm(12), Cm(1.2),
               '개별 솔루션 도입', font_size=13, bold=True, color=C_WHITE, fill_color=RGBColor(0x8B, 0x8B, 0x8B))

rows = [
    ("PII 메타데이터", "1회 스캔 → 전 모듈 공유, 즉시 반영", "솔루션마다 PII 개별 등록, 불일치 위험"),
    ("인프라", "단일 서버, 단일 DB, Docker 1세트", "솔루션별 서버/DB 별도, 3~4배 비용"),
    ("운영 관리", "통합 관리 콘솔 1개, 학습 1회", "솔루션별 관리 화면 별도, 제품별 교육"),
    ("대시보드/보고", "통합 CPO 대시보드, 한눈에 파악", "솔루션별 리포트 수동 취합"),
    ("도입 기간", "4모듈 동시 구축, 공통 기반 활용", "순차 도입 → 연동 공수 추가"),
    ("비용", "통합 라이선스, 유지보수 1건", "개별 라이선스 ×3~4, 유지보수 ×3~4"),
    ("감사 대응", "\"X-One 도입\" 1건으로 전 영역 증빙", "영역별 솔루션 증빙 개별 제출"),
]

for i, (label, xone, others) in enumerate(rows):
    y = Cm(4.4 + i * 1.65)
    bg = C_VERY_LIGHT if i % 2 == 0 else C_WHITE
    add_shape(slide, Cm(1.5), y, Cm(31), Cm(1.65), fill_color=bg)
    add_shape_text(slide, Cm(1.5), y, Cm(7), Cm(1.65),
                   label, font_size=12, bold=True, color=C_WHITE, fill_color=C_NAVY if i % 2 == 0 else C_DEEP_BLUE)
    add_text(slide, Cm(9), y, Cm(11.5), Cm(1.65),
             xone, font_size=11, color=C_TEAL, valign='middle')
    add_text(slide, Cm(21), y, Cm(11.5), Cm(1.65),
             others, font_size=11, color=C_GRAY, valign='middle')

add_key_point(slide, 'PII 메타데이터 공유가 All-in-One의 근본 경쟁력 → X-Scan이 발견한 PII가 파기·감시·마스킹에 즉시 활용', top=Cm(16.5))
add_bottom_bar(slide)
add_page_number(slide, 8)


# ═══════════════════════════════════════════════════════
# ═══ Slide 8: 플랫폼 아키텍처 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 7, 'X-One 플랫폼 아키텍처',
                   'Hub 중심 모듈형 설계 — 필요한 모듈만 선택적 도입, 단계적 확장 가능')

# Hub
add_rounded_rect(slide, Cm(9), Cm(3.5), Cm(15.5), Cm(3), fill_color=C_NAVY)
add_text(slide, Cm(9.5), Cm(3.6), Cm(5), Cm(1),
         'X-One Hub', font_size=18, bold=True, color=C_TEAL, valign='middle')
add_text(slide, Cm(9.5), Cm(4.8), Cm(14.5), Cm(1.2),
         '"All Data. One Platform."    통합 대시보드  |  PII 메타 저장소  |  사용자 권한 관리  |  감사 리포트 엔진',
         font_size=10, color=C_WHITE, valign='middle')

# PII 레이어
add_shape_text(slide, Cm(1.5), Cm(7), Cm(30.5), Cm(1),
               '▼  PII 메타데이터 공유 레이어  ▼', font_size=12, bold=True, color=C_WHITE, fill_color=C_TEAL)

# 4 modules
mod_data = [
    ("X-Scan", "개인정보 탐지", "AI + 패턴 + 메타\n자동스캔 · PII 분류", C_ACCENT2),
    ("X-Purge", "개인정보 파기", "No-Code Wizard\n6단계 전자동 Pipeline", C_RED),
    ("X-Audit", "접속기록 · 소명", "5W1H 실시간 수집\n이상탐지 · 소명관리", C_ACCENT),
    ("X-Gen", "테스트데이터 생성", "연관관계 기반\nKeyMap 자동생성", C_TEAL),
]

for i, (name, subtitle, desc, color) in enumerate(mod_data):
    x = Cm(1.5 + i * 8)
    add_rounded_rect(slide, x, Cm(8.8), Cm(7.2), Cm(5.5), fill_color=C_WHITE, line_color=color, line_width=Pt(2))
    add_shape(slide, x, Cm(8.8), Cm(7.2), Cm(0.15), fill_color=color)
    add_text(slide, x, Cm(9.2), Cm(7.2), Cm(1.2),
             name, font_size=17, bold=True, color=color, align=PP_ALIGN.CENTER, valign='middle')
    add_text(slide, x, Cm(10.5), Cm(7.2), Cm(0.7),
             subtitle, font_size=10, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER, valign='middle')
    add_text(slide, Emu(x.emu + Cm(0.5).emu), Cm(11.5), Cm(6.2), Cm(2.3),
             desc, font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER, valign='middle')

# DB layer
add_shape_text(slide, Cm(1.5), Cm(15), Cm(30.5), Cm(1),
               '업무 DB (다중 연결)  —  Oracle  ·  MariaDB  ·  MySQL  ·  PostgreSQL  ·  MS-SQL',
               font_size=10, color=C_BLACK, fill_color=C_LIGHT_GRAY)

add_key_point(slide, '모듈형 구조: 필요한 모듈만 ON/OFF → 단계적 도입 가능    |    PII 메타 저장소: 전 모듈이 공유 → 중복 작업 제거')
add_bottom_bar(slide)
add_page_number(slide, 9)


# ═══════════════════════════════════════════════════════
# ═══ Slide 9: X-Scan ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 8, 'X-Scan — 개인정보 자동 탐지',
                   'AI + 패턴 + 메타 분석으로 전사 DB의 개인정보를 자동 발견·분류')

engines = [
    ("Meta 40%", "컬럼명 · 타입 · 코멘트\n메타 기반 분석",
     ["테이블/컬럼 이름 패턴 분석", "DB 코멘트 자연어 분석", "데이터 타입·길이 매칭"], C_ACCENT2),
    ("Pattern 35%", "정규식 매칭\n데이터 패턴 기반 탐지",
     ["주민번호, 핸드폰, 이메일 등", "한국 특화 패턴 내장", "커스텀 패턴 등록 가능"], C_ACCENT),
    ("AI 25%", "LLM 기반\n의미론적 분석",
     ["문맥 기반 PII 판별", "비정형 데이터 분석", "신규 PII 유형 자동 탐지"], C_PURPLE),
]

for i, (title, subtitle, items, color) in enumerate(engines):
    x = Cm(1.5 + i * 10.8)
    add_rounded_rect(slide, x, Cm(3.5), Cm(10), Cm(11), fill_color=C_WHITE, line_color=color, line_width=Pt(2))

    add_shape_text(slide, x, Cm(3.5), Cm(10), Cm(2),
                   title, font_size=22, bold=True, color=C_WHITE, fill_color=color)
    add_text(slide, x, Cm(5.5), Cm(10), Cm(1.5),
             subtitle, font_size=11, color=color, align=PP_ALIGN.CENTER, valign='middle')

    add_multiline(slide, Emu(x.emu + Cm(1.2).emu), Cm(8), Cm(7.5), Cm(5.5),
                  items, font_size=12, color=C_BLACK)

add_key_point(slide, 'X-Scan의 스캔 결과가 나머지 3개 모듈의 기반 데이터 → 1회 스캔으로 파기·감시·마스킹 대상 자동 설정')
add_bottom_bar(slide)
add_page_number(slide, 10)


# ═══════════════════════════════════════════════════════
# ═══ Slide 10: X-Purge ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 9, 'X-Purge — 개인정보 파기 자동화',
                   '파기의 전 과정을 No-Code Wizard로 자동화, 하루 10분 관리')

features = [
    ("KeyMap 연관 추적", ["이기종 DB 간 관계 자동 추적", "FK/PK 기반 연관 테이블 탐색", "누락 없는 완전파기 보장"]),
    ("No-Code Wizard", ["SQL 없이 클릭만으로 파기 Job 구성", "테이블 선택 → 조인 자동 감지", "실무 담당자 누구나 설정 가능"]),
    ("6단계 Pipeline", ["대상선정 → 백업 → 검증", "→ 삭제 → 확인 → 완료 전자동", "결재 연동 (상신 → 승인)"]),
    ("CPO 대시보드", ["경영진 보고 포맷 자동 생성", "감사 리포트·증빙 자동생성", "파기 현황 실시간 모니터링"]),
]

for i, (title, items) in enumerate(features):
    x = Cm(1.5 + i * 8)
    add_rounded_rect(slide, x, Cm(3.5), Cm(7.2), Cm(10.5), fill_color=C_WHITE, line_color=C_RED, line_width=Pt(1))
    add_shape_text(slide, x, Cm(3.5), Cm(7.2), Cm(1.5),
                   title, font_size=14, bold=True, color=C_WHITE, fill_color=C_RED)
    add_multiline(slide, Emu(x.emu + Cm(0.8).emu), Cm(5.8), Cm(5.8), Cm(7),
                  items, font_size=11, color=C_BLACK)

add_text(slide, Cm(1.5), Cm(14.5), Cm(31), Cm(0.8),
         '프로세스:  Discovery → Policy → Job → 결재 → 실행(스케줄) → 모니터링 → 증빙',
         font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER, valign='middle')

add_key_point(slide, '초기 설정 1회 후, 매월 전자동 실행. 담당자는 대시보드 확인만 하면 됩니다.')
add_bottom_bar(slide)
add_page_number(slide, 11)


# ═══════════════════════════════════════════════════════
# ═══ Slide 11: No-Code Wizard 비교 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 10, 'X-Purge No-Code Wizard — SQL 없이 클릭만으로',
                   '초기 설정 1회, 이후 전자동 실행')

# 기존
add_shape_text(slide, Cm(1.5), Cm(3.2), Cm(15.5), Cm(1.3),
               '기존: SQL 직접 입력 방식', font_size=14, bold=True, color=C_WHITE, fill_color=C_GRAY)
old_steps = [
    "1.  파기 대상 테이블 분석 (ERD 확인, 관계 파악)",
    "2.  연관 테이블 조인 SQL 작성 (수십 줄 쿼리)",
    "3.  테스트 환경에서 SQL 검증 (오류 반복)",
    "4.  운영 반영 및 수동 실행 (DBA 필수)",
    "5.  결과 확인 SQL 별도 작성",
    "6.  매번 변경 시 1~5 반복",
]
for j, step in enumerate(old_steps):
    add_text(slide, Cm(2), Cm(4.8 + j * 1.5), Cm(14.5), Cm(1.2),
             step, font_size=11, color=C_RED, valign='middle')
add_text(slide, Cm(2), Cm(14.2), Cm(14.5), Cm(0.8),
         'SQL 역량 필수  |  전문 DBA 필요  |  퇴사 시 유실', font_size=11, bold=True, color=C_RED)

# X-Purge
add_shape_text(slide, Cm(17.5), Cm(3.2), Cm(15.5), Cm(1.3),
               'X-Purge: No-Code Wizard 방식', font_size=14, bold=True, color=C_WHITE, fill_color=C_TEAL)
new_steps = [
    "1.  테이블 검색 → 클릭 선택 (1분)",
    "2.  조인 컬럼 자동 감지 → 확인 클릭",
    "3.  파기 조건 Wizard로 설정 (날짜/조건)",
    "4.  저장 → 결재 → 스케줄 자동 실행",
    "5.  대시보드에서 결과 자동 확인",
    "6.  변경 시 Wizard 재설정 (1분)",
]
for j, step in enumerate(new_steps):
    add_text(slide, Cm(18), Cm(4.8 + j * 1.5), Cm(14.5), Cm(1.2),
             step, font_size=11, color=C_TEAL, valign='middle')
add_text(slide, Cm(18), Cm(14.2), Cm(14.5), Cm(0.8),
         'SQL 지식 불필요  |  실무 담당자 누구나 설정 가능', font_size=11, bold=True, color=C_TEAL)

add_key_point(slide, 'DBA 의존 제거 → 담당자 누구나 파기 설정·실행·모니터링 가능')
add_bottom_bar(slide)
add_page_number(slide, 12)


# ═══════════════════════════════════════════════════════
# ═══ X-Audit 섹션 (Slide 12 ~ 18) ═══
# ═══════════════════════════════════════════════════════

# ── Slide 12: X-Audit 전체 개요 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide, C_ACCENT)
add_section_header(slide, 11, 'X-Audit — 접속기록 수집 · 이상행위 탐지 · 소명',
                   '개인정보 접근의 5W1H를 실시간 수집하고, 이상행위를 자동 탐지하고, 소명까지 완결합니다',
                   header_color=C_ACCENT)

caps = [
    ("5W1H 수집", "접속기록\n실시간 수집", "DB Audit + WAS Log\n+ Java Agent", C_ACCENT2),
    ("이상행위 탐지", "7가지 규칙\n자동 탐지", "대량조회/야간/PII등급\n미등록IP/반복접근", C_RED),
    ("소명 관리", "자동 소명요청\n→ 검토 → 승인", "기한관리/결과기록\nAPI 자동등록", C_ACCENT),
    ("위변조 방지", "SHA-256\n해시 체인", "레코드별 연쇄해시\n변조 즉시 감지", C_DEEP_BLUE),
    ("감사 보고", "월간점검\n자동 생성", "CPO 보고용 포맷\n감사 증빙 완결", C_TEAL),
]

for i, (title, sub, desc, color) in enumerate(caps):
    x = Cm(1 + i * 6.4)
    add_rounded_rect(slide, x, Cm(3.5), Cm(5.8), Cm(10.5), fill_color=C_WHITE, line_color=color, line_width=Pt(2))
    add_shape_text(slide, x, Cm(3.5), Cm(5.8), Cm(1.5),
                   title, font_size=14, bold=True, color=C_WHITE, fill_color=color)
    add_text(slide, Emu(x.emu + Cm(0.3).emu), Cm(5.8), Cm(5.2), Cm(3),
             sub, font_size=13, bold=True, color=color, align=PP_ALIGN.CENTER, valign='middle')
    add_text(slide, Emu(x.emu + Cm(0.3).emu), Cm(9.5), Cm(5.2), Cm(3),
             desc, font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER, valign='middle')

add_key_point(slide, 'X-Scan PII 메타 연동: 접속 시 자동으로 PII 등급 분류 → 1급 개인정보 대량 접근 시 즉시 알림 + 소명 요청', top=Cm(14.8))
add_text(slide, Cm(1.5), Cm(16), Cm(31), Cm(0.8),
         '개인정보보호법 안전성 확보조치 제8조: 접속기록 최소 1년(금융사 5년) | 월 1회 점검 | 위변조 방지 의무',
         font_size=9, color=C_GRAY, align=PP_ALIGN.CENTER, valign='middle')
add_bottom_bar(slide, C_ACCENT)
add_page_number(slide, 13)


# ── Slide 13: 5W1H 수집 상세 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide, C_ACCENT)
add_section_header(slide, 12, 'X-Audit 접속기록 수집 — 5W1H 완전 기록',
                   header_color=RGBColor(0x8B, 0x4D, 0x00))

w5h1 = [
    ("Who", "접속자 계정", "DB 계정 + WAS 실사용자 매핑\n공용계정도 실사용자 추적 가능"),
    ("When", "접속 일시", "밀리초 단위 정밀 기록\n시간대별 패턴 분석 기반"),
    ("Where", "접속 IP / 단말", "IP, 호스트명, 접속 경로\nVPN/NAT 환경 대응"),
    ("What", "작업 유형", "SELECT/INSERT/UPDATE/DELETE\nDDL(DROP, ALTER) 별도 분류"),
    ("Whom", "정보주체 식별", "조회 대상 고객 식별\n대량 조회 시 건수 기록"),
    ("How", "접근 경로", "화면ID, API 경로, 배치Job명\nJava Agent 실시간 캡처"),
]

for i, (label, title, desc) in enumerate(w5h1):
    row = i // 3
    col = i % 3
    x = Cm(1.5 + col * 10.8)
    y = Cm(3 + row * 6.8)

    add_rounded_rect(slide, x, y, Cm(10), Cm(5.8), fill_color=C_WHITE, line_color=C_ACCENT, line_width=Pt(1.5))

    add_shape_text(slide, Emu(x.emu + Cm(0.5).emu), Emu(y.emu + Cm(0.8).emu), Cm(2.5), Cm(2.5),
                   label, font_size=16, bold=True, color=C_WHITE, fill_color=C_ACCENT, shape_type='oval')

    add_text(slide, Emu(x.emu + Cm(3.5).emu), Emu(y.emu + Cm(0.5).emu), Cm(6), Cm(1.2),
             title, font_size=14, bold=True, color=C_BLACK, valign='middle')
    add_text(slide, Emu(x.emu + Cm(3.5).emu), Emu(y.emu + Cm(2.2).emu), Cm(6), Cm(3),
             desc, font_size=11, color=C_GRAY, valign='top')

add_key_point(slide, '법적 요구사항인 "누가, 언제, 어디서, 무엇을, 어떻게" 접근했는지를 자동으로 완전 기록합니다.')
add_bottom_bar(slide, C_ACCENT)
add_page_number(slide, 14)


# ── Slide 14: 로그 수집 전략 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide, C_ACCENT)
add_section_header(slide, 13, 'X-Audit 로그 수집 전략 — 리스크 기반 선별 로깅',
                   header_color=RGBColor(0x8B, 0x4D, 0x00))

add_text(slide, Cm(1.5), Cm(2.8), Cm(31), Cm(0.8),
         '"모든 것을 감시하는 것은 아무것도 감시하지 않는 것과 같습니다"',
         font_size=13, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER, valign='middle')

steps_data = [
    ("STEP 1", "Selection", "데이터 등급 필터", ["S급: 주민번호, 외국인번호", "A급: 계좌, 연락처 집중", "B급: 성명 등 → 통계만"]),
    ("STEP 2", "Grading", "접근 주체별 로그 레벨", ["관리자/개발자 → Full SQL", "배치/시스템 → 통계만", "WAS 서비스 → 패턴 로깅"]),
    ("STEP 3", "Thresholding", "임계치 기반 동적 로깅", ["대량 조회 시 SQL 강제 저장", "야간/휴일 자동 격상", "Threshold 초과 시 상세기록"]),
    ("STEP 4", "Flow Mapping", "길목(Choke Point) 전략", ["통합 원장/결과 마트 집중", "대외기관 송수신 테이블", "무조건 감시 대상 등록"]),
]

for i, (step, en, title, items) in enumerate(steps_data):
    x = Cm(1.5 + i * 8)
    add_rounded_rect(slide, x, Cm(4), Cm(7.2), Cm(8.5), fill_color=C_WHITE, line_color=C_ACCENT, line_width=Pt(1.5))
    add_shape_text(slide, x, Cm(4), Cm(7.2), Cm(1),
                   step, font_size=10, color=C_WHITE, fill_color=C_ACCENT)
    add_text(slide, x, Cm(5), Cm(7.2), Cm(1),
             en, font_size=15, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER, valign='middle')
    add_text(slide, Emu(x.emu + Cm(0.5).emu), Cm(6.2), Cm(6.2), Cm(0.8),
             title, font_size=11, bold=True, color=C_BLACK, valign='middle')
    add_multiline(slide, Emu(x.emu + Cm(0.5).emu), Cm(7.3), Cm(6.2), Cm(4.5),
                  items, font_size=10, color=C_GRAY)

# 수집 방식
add_text(slide, Cm(1.5), Cm(13), Cm(31), Cm(0.7),
         '수집 방식별 역할', font_size=13, bold=True, color=C_BLACK)

methods = [
    ("DB Audit Log", "S/A급 주요 개인정보 테이블 선별 수집\nDB 레벨 SQL 원본 확보"),
    ("WAS Log 기반", "사용자·화면·업무 컨텍스트 확보\nDB Audit만으로 알 수 없는 Who/Why 보강"),
    ("Java Agent", "JDBC 레벨 실시간 가로채기\n실시간 SQL + 사용자 식별 + 응답시간"),
]
for i, (name, desc) in enumerate(methods):
    x = Cm(1.5 + i * 10.8)
    add_rounded_rect(slide, x, Cm(13.8), Cm(10), Cm(2.2), fill_color=C_VERY_LIGHT, line_color=C_ACCENT, line_width=Pt(1))
    add_text(slide, Emu(x.emu + Cm(0.5).emu), Cm(13.9), Cm(2.8), Cm(0.7),
             name, font_size=10, bold=True, color=C_ACCENT, valign='middle')
    add_text(slide, Emu(x.emu + Cm(3.5).emu), Cm(13.9), Cm(6.2), Cm(2),
             desc, font_size=9, color=C_GRAY, valign='middle')

add_key_point(slide, '기대 효과: 연산 부하 60% 감소  |  저장 용량 80% 절감  |  유의미 사고 데이터 집중  |  시스템 성능 유지', top=Cm(16.5))
add_bottom_bar(slide, C_ACCENT)
add_page_number(slide, 15)


# ── Slide 15: 이상행위 탐지 7가지 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide, C_ACCENT)
add_section_header(slide, 14, 'X-Audit 이상행위 탐지 — 7가지 규칙 엔진',
                   header_color=RGBColor(0x8B, 0x4D, 0x00))

rules = [
    ("VOLUME", "대량 접속 감지", "시간 윈도우 내 대량 조회/다운로드\n임계값 초과 시 즉시 알림", "CRITICAL"),
    ("TIME_RANGE", "야간/휴일 접근", "정상 업무시간 외 접근 자동 탐지\n공휴일/특수일 자동 반영", "HIGH"),
    ("ACCESS_DENIED", "접근 거부 반복", "로그인 실패, 권한없는 접근 반복\n브루트포스 공격 패턴 감지", "HIGH"),
    ("PII_GRADE", "고등급 PII 접근", "S/A급 개인정보 대량 접근 감지\nX-Scan PII 메타 자동 연동", "CRITICAL"),
    ("REPEAT", "동일 테이블 반복", "동일 테이블 비정상 반복 접근\n데이터 유출 시도 패턴 감지", "MEDIUM"),
    ("NEW_IP", "미등록 IP 접근", "90일 이력 기반 신규 IP 탐지\nVPN/프록시 우회 접근 감지", "HIGH"),
    ("INACTIVE", "미사용 계정", "장기 미사용 계정 접근 탐지\n퇴직자/비활성 계정 감시", "MEDIUM"),
]

for i, (code, title, desc, severity) in enumerate(rules):
    row = i // 4
    col = i % 4
    x = Cm(1 + col * 8.2)
    y = Cm(3 + row * 7.2)

    sev_color = C_RED if severity == "CRITICAL" else C_ACCENT if severity == "HIGH" else C_ACCENT2

    add_rounded_rect(slide, x, y, Cm(7.5), Cm(6.2), fill_color=C_WHITE, line_color=sev_color, line_width=Pt(1.5))

    add_shape_text(slide, Emu(x.emu + Cm(0.4).emu), Emu(y.emu + Cm(0.4).emu), Cm(4), Cm(0.9),
                   code, font_size=10, bold=True, color=C_WHITE, fill_color=sev_color, shape_type='rounded')
    add_text(slide, Emu(x.emu + Cm(4.8).emu), Emu(y.emu + Cm(0.4).emu), Cm(2.5), Cm(0.9),
             severity, font_size=9, bold=True, color=sev_color, align=PP_ALIGN.RIGHT, valign='middle')

    add_text(slide, Emu(x.emu + Cm(0.5).emu), Emu(y.emu + Cm(1.8).emu), Cm(6.5), Cm(1),
             title, font_size=13, bold=True, color=C_BLACK, valign='middle')
    add_text(slide, Emu(x.emu + Cm(0.5).emu), Emu(y.emu + Cm(3.2).emu), Cm(6.5), Cm(2.5),
             desc, font_size=10, color=C_GRAY, valign='top')

add_key_point(slide, '7가지 규칙은 독립 실행되며 복합 조건 조합 가능 → 탐지 후 자동으로 소명 프로세스 시작')
add_bottom_bar(slide, C_ACCENT)
add_page_number(slide, 16)


# ── Slide 16: 소명 프로세스 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide, C_ACCENT)
add_section_header(slide, 15, 'X-Audit 소명 프로세스 — 탐지에서 감사까지 완결',
                   header_color=RGBColor(0x8B, 0x4D, 0x00))

add_text(slide, Cm(1.5), Cm(2.9), Cm(31), Cm(0.7),
         '개인정보보호법 안전성 확보조치 제8조 2항 — "이상행위 확인 시 소명 절차 수행"',
         font_size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# 워크플로우
workflow = [
    ("이상행위\n자동 탐지", "실시간\n7가지 규칙", C_RED),
    ("소명 요청\n자동 발송", "이메일 +\n대시보드", C_ACCENT),
    ("담당자\n소명 작성", "사유 설명\n기한 관리", C_ACCENT2),
    ("관리자\n검토/승인", "승인/반려\n의견 기록", C_DEEP_BLUE),
    ("보고서\n자동 반영", "월간점검\n감사 증빙", C_TEAL),
]

for i, (title, desc, color) in enumerate(workflow):
    x = Cm(1.2 + i * 6.4)
    add_rounded_rect(slide, x, Cm(4), Cm(5.8), Cm(5.5), fill_color=C_WHITE, line_color=color, line_width=Pt(2))

    add_shape_text(slide, Emu(x.emu + Cm(1.9).emu), Cm(4.5), Cm(2), Cm(2),
                   str(i+1), font_size=18, bold=True, color=C_WHITE, fill_color=color, shape_type='oval')

    add_text(slide, x, Cm(6.8), Cm(5.8), Cm(1.8),
             title, font_size=12, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER, valign='middle')
    add_text(slide, x, Cm(8.2), Cm(5.8), Cm(1.2),
             desc, font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER, valign='middle')

    if i < 4:
        add_text(slide, Cm(6.5 + i * 6.4), Cm(6.5), Cm(1.2), Cm(1.5),
                 '→', font_size=24, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER, valign='middle')

# 상세 기능
details = [
    ("자동 소명 요청", "이상행위 탐지 즉시 담당자에게\n이메일/대시보드 알림 자동 발송\n소명 기한 자동 설정 및 독촉"),
    ("소명 작성 지원", "이상행위 상세 내역 자동 첨부\n사유 템플릿 제공 (업무/장애 등)\n첨부파일 업로드 지원"),
    ("검토/승인 관리", "관리자 검토 → 승인/반려 판정\n반려 시 재소명 요청 자동 발송\n전체 이력 완전 보존"),
    ("감사 증빙 완결", "월간점검 보고서 자동 포함\nCPO 보고용 포맷 자동 생성\n외부 감사 대응 원클릭 출력"),
]

for i, (title, desc) in enumerate(details):
    x = Cm(1.5 + i * 8)
    add_rounded_rect(slide, x, Cm(10.5), Cm(7.2), Cm(5), fill_color=C_VERY_LIGHT, line_color=C_ACCENT, line_width=Pt(1))
    add_text(slide, Emu(x.emu + Cm(0.5).emu), Cm(10.8), Cm(6.2), Cm(0.8),
             title, font_size=12, bold=True, color=C_ACCENT, valign='middle')
    add_text(slide, Emu(x.emu + Cm(0.5).emu), Cm(12), Cm(6.2), Cm(3),
             desc, font_size=10, color=C_GRAY, valign='top')

add_key_point(slide, '소명 결과 API 자동 등록: 탐지 → 소명요청 → 응답 → API 저장 → 보고서 자동 포함 → 감사 대응 완결')
add_bottom_bar(slide, C_ACCENT)
add_page_number(slide, 17)


# ── Slide 17: 위변조/보관/보고서 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide, C_ACCENT)
add_section_header(slide, 16, 'X-Audit 위변조 방지 · 보관 · 월간점검 보고서',
                   header_color=RGBColor(0x8B, 0x4D, 0x00))

pillars = [
    ("위변조 방지", C_RED,
     ["SHA-256 해시 체인 적용", "레코드별 연쇄해시 생성", "변조 시 즉시 감지", "무결성 검증 API 제공", "감사 시 원클릭 검증"]),
    ("자동 보관/아카이브", C_DEEP_BLUE,
     ["월 단위 파티션 자동 생성", "금융사 5년 보관 준수", "보관 기간 만료 자동 삭제", "아카이브 자동 관리", "저장소 용량 최적화"]),
    ("월간점검 보고서", C_TEAL,
     ["접근패턴 분석 리포트", "이상행위·소명 처리 현황", "CPO 보고용 자동 생성", "기간별/부서별 통계", "외부 감사 대응 포맷"]),
]

for i, (title, color, items) in enumerate(pillars):
    x = Cm(1.5 + i * 10.8)
    add_rounded_rect(slide, x, Cm(3.5), Cm(10), Cm(12), fill_color=C_WHITE, line_color=color, line_width=Pt(2))
    add_shape_text(slide, x, Cm(3.5), Cm(10), Cm(2.2),
                   title, font_size=20, bold=True, color=C_WHITE, fill_color=color)
    add_multiline(slide, Emu(x.emu + Cm(1.2).emu), Cm(6.8), Cm(7.5), Cm(8),
                  items, font_size=12, color=C_BLACK, spacing=1.6)

add_key_point(slide, '접속기록의 저장 → 보관 → 점검 → 보고서까지 전 주기를 자동화하여 안전성 확보조치를 완벽 이행합니다.')
add_bottom_bar(slide, C_ACCENT)
add_page_number(slide, 18)


# ── Slide 18: X-Audit 경쟁 우위 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide, C_ACCENT)
add_section_header(slide, 17, 'X-Audit vs 경쟁 솔루션 — 압도적 차별화',
                   header_color=RGBColor(0x8B, 0x4D, 0x00))

# Header row
add_shape_text(slide, Cm(1.5), Cm(3.2), Cm(5.5), Cm(1.3),
               '구분', font_size=12, bold=True, color=C_WHITE, fill_color=C_GRAY)
add_shape_text(slide, Cm(7), Cm(3.2), Cm(12.5), Cm(1.3),
               'X-Audit', font_size=12, bold=True, color=C_WHITE, fill_color=C_ACCENT)
add_shape_text(slide, Cm(19.5), Cm(3.2), Cm(13), Cm(1.3),
               '경쟁 솔루션', font_size=12, bold=True, color=C_WHITE, fill_color=RGBColor(0x8B, 0x8B, 0x8B))

comp_rows = [
    ("수집 방식", "DB Audit + WAS Log + Java Agent\n3중 수집으로 5W1H 완전 확보", "DB Audit 단일 또는 전용 에이전트\n사용자/화면 컨텍스트 부재"),
    ("이상행위 탐지", "7가지 규칙 엔진 + PII 등급 연동\n복합 조건 조합 가능", "기본 규칙 2~3개\nPII 등급 미연동"),
    ("소명 관리", "자동 소명요청 → 작성 → 검토 → 승인\nAPI 기반 전 과정 자동화", "소명 기능 미지원 또는\n수동 엑셀 관리"),
    ("위변조 방지", "SHA-256 해시 체인\n레코드별 연쇄해시, 검증 API", "미지원 또는 제한적\n별도 솔루션 필요"),
    ("PII 메타 연동", "X-Scan 자동 연동\n접근 시 PII 등급 자동 분류", "PII 정보 수동 등록\n또는 미연동"),
    ("보고서", "월간점검 보고서 자동 생성\nCPO 대시보드 통합", "수동 리포트 작성\n별도 BI 도구 필요"),
]

for i, (label, xaudit, others) in enumerate(comp_rows):
    y = Cm(4.5 + i * 2.1)
    bg = C_VERY_LIGHT if i % 2 == 0 else C_WHITE
    add_shape(slide, Cm(1.5), y, Cm(31), Cm(2.1), fill_color=bg)
    add_shape_text(slide, Cm(1.5), y, Cm(5.5), Cm(2.1),
                   label, font_size=11, bold=True, color=C_WHITE, fill_color=C_NAVY if i % 2 == 0 else C_DEEP_BLUE)
    add_text(slide, Cm(7.5), y, Cm(11.5), Cm(2.1),
             xaudit, font_size=10, color=C_TEAL, valign='middle')
    add_text(slide, Cm(20), y, Cm(12.5), Cm(2.1),
             others, font_size=10, color=C_GRAY, valign='middle')

add_key_point(slide, 'X-Audit은 수집·탐지·소명·보고의 전 과정을 자동화한 유일한 접속기록관리 솔루션입니다.', top=Cm(17))
add_bottom_bar(slide, C_ACCENT)
add_page_number(slide, 19)


# ═══════════════════════════════════════════════════════
# ═══ Slide 19: X-Gen ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 18, 'X-Gen — 테스트데이터 자동생성 (도입 완료)',
                   '운영 DB 연관관계 기반으로 안전한 테스트데이터를 자동 생성')

features = [
    ("연관관계 자동추적", ["운영 DB FK/PK 관계 자동 분석", "다중 테이블 관계 무결성 보장", "이기종 DB 간 관계도 추적"], C_TEAL),
    ("PII 자동 마스킹", ["X-Scan PII 메타 연동 자동 변환", "SQL 작성 불필요 — 원클릭 마스킹", "1:1 비율 보장 (통계적 동치)"], C_ACCENT2),
    ("신청-승인 워크플로", ["개발자 웹에서 데이터 신청", "관리자 승인 후 자동 생성", "결재 기반 거버넌스 통제"], C_ACCENT),
    ("사용 현황 관리", ["누가/언제/얼마나 사용 추적", "리포트 자동 생성", "테스트 데이터 라이프사이클 관리"], C_GREEN),
]

for i, (title, items, color) in enumerate(features):
    x = Cm(1.5 + i * 8)
    add_rounded_rect(slide, x, Cm(3.5), Cm(7.2), Cm(10), fill_color=C_WHITE, line_color=color, line_width=Pt(1.5))
    add_shape_text(slide, x, Cm(3.5), Cm(7.2), Cm(1.5),
                   title, font_size=13, bold=True, color=C_WHITE, fill_color=color)
    add_multiline(slide, Emu(x.emu + Cm(0.8).emu), Cm(5.8), Cm(5.8), Cm(7),
                  items, font_size=11, color=C_BLACK)

add_text(slide, Cm(1.5), Cm(14), Cm(31), Cm(0.8),
         'JB우리캐피탈은 X-Gen이 이미 도입되어 있으므로 KeyMap 데이터 확보 → X-Purge 즉시 적용 가능',
         font_size=12, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER, valign='middle')

add_key_point(slide, 'X-Gen이 구축한 연관관계(KeyMap) 데이터는 X-Purge의 파기 대상 추적에 그대로 재활용됩니다.')
add_bottom_bar(slide)
add_page_number(slide, 20)


# ═══════════════════════════════════════════════════════
# ═══ Slide 20: 모듈 간 시너지 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 19, 'X-One 모듈 간 데이터 시너지',
                   '4개 모듈이 PII 메타와 연관관계를 공유하여 All-in-One 가치를 극대화')

add_shape_text(slide, Cm(8), Cm(3.5), Cm(17.5), Cm(2),
               'X-Scan — PII 자동 탐지     테이블A.컬럼1 = 주민번호(S급)   테이블B.컬럼3 = 전화번호(A급)',
               font_size=11, bold=True, color=C_WHITE, fill_color=C_ACCENT2, shape_type='rounded')

add_text(slide, Cm(14.5), Cm(5.5), Cm(4.5), Cm(1),
         '▼  PII 메타 자동 공유  ▼', font_size=11, bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)

receivers = [
    ("X-Purge", "테이블A에 S급 PII 있음\n→ 파기 대상 자동 포함", C_RED),
    ("X-Audit", "테이블A 접근 시 S급 PII 접근으로\n자동 분류 → 즉시 알림 + 소명", C_ACCENT),
    ("X-Gen", "테이블A 복사 시 주민번호/이름\n자동 마스킹 → SQL 없이 자동 변환", C_TEAL),
]
for i, (name, desc, color) in enumerate(receivers):
    x = Cm(1.5 + i * 10.8)
    add_rounded_rect(slide, x, Cm(7), Cm(10), Cm(4.5), fill_color=C_WHITE, line_color=color, line_width=Pt(2))
    add_text(slide, x, Cm(7.3), Cm(10), Cm(1),
             name, font_size=15, bold=True, color=color, align=PP_ALIGN.CENTER, valign='middle')
    add_text(slide, Emu(x.emu + Cm(0.5).emu), Cm(8.5), Cm(9), Cm(2.5),
             desc, font_size=11, color=C_BLACK, align=PP_ALIGN.CENTER, valign='middle')

add_shape_text(slide, Cm(1.5), Cm(12.5), Cm(31), Cm(1.5),
               'X-Gen ──── KeyMap 데이터 (테이블 간 관계) ────→ X-Purge (연관 파기)',
               font_size=11, color=C_WHITE, fill_color=C_NAVY, shape_type='rounded')

add_key_point(slide, '이것이 개별 솔루션으로는 절대 불가능한 X-One의 핵심 가치입니다.')
add_bottom_bar(slide)
add_page_number(slide, 21)


# ═══════════════════════════════════════════════════════
# ═══ Slide 21: 경쟁 비교 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 20, '경쟁 솔루션과의 차이',
                   'X-One은 통합 플랫폼 설계로 근본적으로 다릅니다')

add_shape_text(slide, Cm(1.5), Cm(3.2), Cm(5.5), Cm(1.2),
               '구분', font_size=12, bold=True, color=C_WHITE, fill_color=C_GRAY)
add_shape_text(slide, Cm(7), Cm(3.2), Cm(12.5), Cm(1.2),
               'X-One (All-in-One)', font_size=12, bold=True, color=C_WHITE, fill_color=C_TEAL)
add_shape_text(slide, Cm(19.5), Cm(3.2), Cm(13), Cm(1.2),
               '경쟁사', font_size=12, bold=True, color=C_WHITE, fill_color=RGBColor(0x8B, 0x8B, 0x8B))

comp_data = [
    ("제품 구조", "1 플랫폼 · 4 모듈", "영역별 별도 제품"),
    ("PII 관리", "1회 스캔 → 전 모듈 공유", "제품마다 개별 PII 등록"),
    ("파기 자동화", "전자동 No-Code Wizard", "반자동, SQL 직접 구성"),
    ("연관 데이터 추적", "KeyMap 자동 추적", "미지원 또는 제한적"),
    ("접속기록 수집", "DB Audit + WAS Log + Agent", "DB Audit만 또는 전용 에이전트"),
    ("이상행위 + 소명", "7규칙 탐지 + 소명 API 자동화", "기본 규칙, 소명 미지원"),
    ("위변조 방지", "SHA-256 해시 체인", "미지원 또는 제한적"),
    ("테스트데이터", "연관관계 기반 자동생성", "별도 솔루션 필요"),
    ("감사 보고", "통합 CPO 대시보드", "영역별 수동 취합"),
]

for i, (label, xone, others) in enumerate(comp_data):
    y = Cm(4.4 + i * 1.35)
    bg = C_VERY_LIGHT if i % 2 == 0 else C_WHITE
    add_shape(slide, Cm(1.5), y, Cm(31), Cm(1.35), fill_color=bg)
    add_shape_text(slide, Cm(1.5), y, Cm(5.5), Cm(1.35),
                   label, font_size=10, bold=True, color=C_WHITE, fill_color=C_NAVY if i % 2 == 0 else C_DEEP_BLUE)
    add_text(slide, Cm(7.5), y, Cm(12), Cm(1.35),
             xone, font_size=10, color=C_TEAL, valign='middle')
    add_text(slide, Cm(20), y, Cm(12.5), Cm(1.35),
             others, font_size=10, color=C_GRAY, valign='middle')

add_bottom_bar(slide)
add_page_number(slide, 22)


# ═══════════════════════════════════════════════════════
# ═══ Slide 22: 도입 효과 ROI ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 21, '도입 효과 & ROI',
                   '측정 가능한 비즈니스 임팩트')

roi_items = [
    ("93%↓", "운영 공수 절감", "파기: 월 40시간 → 3.5시간\n접근이력: 수작업 → 4.5시간/월\n전체 담당자 하루 20분 미만", C_TEAL),
    ("95%↓", "점검 시간 단축", "감사 리포트 자동 생성\n월간점검 보고서 자동 생성\n통합 CPO 대시보드 즉시 보고", C_ACCENT2),
    ("100%", "파기 완전성 확보", "KeyMap 연관 추적\n물리삭제 + 분리보관 + 영구파기\n복원·열람 지원", C_GREEN),
    ("ZERO", "접근이력 사각지대", "5W1H 전수 수집\n이상행위 실시간 탐지\n소명 체계로 사후 관리 완결", C_ACCENT),
]

for i, (big_num, title, desc, color) in enumerate(roi_items):
    x = Cm(1.5 + i * 8)
    add_rounded_rect(slide, x, Cm(3.5), Cm(7.2), Cm(9.5), fill_color=C_WHITE, line_color=color, line_width=Pt(2))

    add_text(slide, x, Cm(4), Cm(7.2), Cm(2.5),
             big_num, font_size=34, bold=True, color=color, align=PP_ALIGN.CENTER, valign='middle')
    add_text(slide, x, Cm(6.5), Cm(7.2), Cm(1),
             title, font_size=14, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER, valign='middle')
    add_shape(slide, Emu(x.emu + Cm(1.5).emu), Cm(7.8), Cm(4.2), Cm(0.05), fill_color=C_LIGHT_GRAY)
    add_text(slide, Emu(x.emu + Cm(0.5).emu), Cm(8.2), Cm(6.2), Cm(4.5),
             desc, font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER, valign='top')

# 추가 가치
add_rounded_rect(slide, Cm(1.5), Cm(13.8), Cm(31), Cm(2.5), fill_color=C_NAVY)
add_text(slide, Cm(2.5), Cm(14), Cm(30), Cm(0.8),
         '추가 가치: 이중 과징금 감경 수단', font_size=14, bold=True, color=C_GOLD)
add_text(slide, Cm(2.5), Cm(15), Cm(30), Cm(1),
         '개정법은 파기 솔루션 + 접속기록 관리 솔루션 투자를 모두 감경 사유로 인정. X-One 하나의 도입으로 4개 영역의 투자 증빙을 한꺼번에 확보합니다.',
         font_size=11, color=C_WHITE, valign='middle')

add_bottom_bar(slide)
add_page_number(slide, 23)


# ═══════════════════════════════════════════════════════
# ═══ Slide 23: 월간 운영 업무량 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 22, 'X-One 도입 후 월간 운영 업무량',
                   '기준: 초기 설정 완료 후 정상 운영')

# Header
cols = [("모듈", Cm(4)), ("월간 운영 업무 항목", Cm(13.5)), ("담당자", Cm(4)), ("시간", Cm(3)), ("비고", Cm(6.5))]
cx = Cm(1.5)
for name, width in cols:
    add_shape_text(slide, cx, Cm(3.2), width, Cm(1), name, font_size=10, bold=True, color=C_WHITE, fill_color=C_NAVY)
    cx = Emu(cx.emu + width.emu)

ops_data = [
    ("X-Scan", "신규 테이블/컬럼 PII 확인·처리", "정보보호", "1.5h", "주당 20분"),
    ("X-Scan", "PII 메타 확정 (개인정보 항목 지정)", "IT", "1.0h", "Discovery 참조"),
    ("X-Purge", "신규 테이블 Job/Step 추가", "IT", "1.0h", "Wizard 1분/건"),
    ("X-Purge", "월간 파기 실행 모니터링", "IT", "1.0h", "대시보드 확인"),
    ("X-Purge", "오류 발생 시 확인·재실행", "IT", "0.5h", "원클릭 재실행"),
    ("X-Purge", "감사 리포트·증빙 확인", "정보보호", "0.5h", "자동 생성"),
    ("X-Audit", "이상행위 알림 확인 & 소명 검토", "정보보호", "2.0h", "알림 시 수시"),
    ("X-Audit", "월간 접속기록 점검", "정보보호", "1.0h", "자동분석 리뷰"),
    ("X-Audit", "수집 상태 모니터링", "IT", "0.5h", "대시보드 확인"),
    ("X-Audit", "월간점검 보고서 확인/제출", "정보보호", "0.5h", "자동 생성"),
]

module_colors = {"X-Scan": C_ACCENT2, "X-Purge": C_RED, "X-Audit": C_ACCENT}

for i, (module, task, person, hours, note) in enumerate(ops_data):
    y = Cm(4.2 + i * 1.1)
    bg = C_VERY_LIGHT if i % 2 == 0 else C_WHITE
    add_shape(slide, Cm(1.5), y, Cm(31), Cm(1.1), fill_color=bg)

    m_color = module_colors.get(module, C_GRAY)
    add_shape_text(slide, Cm(1.7), Emu(y.emu + Cm(0.1).emu), Cm(3.6), Cm(0.9),
                   module, font_size=9, bold=True, color=C_WHITE, fill_color=m_color, shape_type='rounded')
    add_text(slide, Cm(5.5), y, Cm(13.5), Cm(1.1), task, font_size=10, color=C_BLACK, valign='middle')
    add_text(slide, Cm(19), y, Cm(4), Cm(1.1), person, font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER, valign='middle')
    add_text(slide, Cm(23), y, Cm(3), Cm(1.1), hours, font_size=10, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER, valign='middle')
    add_text(slide, Cm(26), y, Cm(6.5), Cm(1.1), note, font_size=10, color=C_GRAY, valign='middle')

add_shape_text(slide, Cm(1.5), Cm(15.3), Cm(31), Cm(1.2),
               '전체 월간 운영 공수: 정보보호 5.5h + IT 4.0h = 총 9.5시간  |  담당자 1인 기준 하루 약 20분 이하',
               font_size=12, bold=True, color=C_WHITE, fill_color=C_NAVY)
add_bottom_bar(slide)
add_page_number(slide, 24)


# ═══════════════════════════════════════════════════════
# ═══ Slide 24: JB 도입 로드맵 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 23, 'JB우리캐피탈 — 도입 로드맵',
                   'X-Gen 기 도입 완료 → X-Purge + X-Audit 추가로 X-One 거버넌스 완성')

# 확보된 자산
add_rounded_rect(slide, Cm(1.5), Cm(3.2), Cm(15.5), Cm(5.5), fill_color=C_VERY_LIGHT, line_color=C_GREEN, line_width=Pt(2))
add_text(slide, Cm(2), Cm(3.5), Cm(14.5), Cm(0.8),
         'X-Gen 도입으로 이미 확보된 자산', font_size=13, bold=True, color=C_GREEN)
assets = [
    "전사 테이블 연관관계 (KeyMap) 구축 완료 → X-Purge 즉시 활용",
    "PII 포함 테이블/컬럼 메타정보 확보 → X-Scan 기반 확보",
    "이기종 DB 접속 환경 구성 완료 → X-Audit 수집 즉시 연동",
    "운영 환경 인프라 검증 완료 → 추가 인프라 불필요",
]
add_multiline(slide, Cm(2.5), Cm(4.8), Cm(14), Cm(3.5), assets, font_size=11, color=C_BLACK)

# 기간 단축
add_rounded_rect(slide, Cm(17.5), Cm(3.2), Cm(15.5), Cm(5.5), fill_color=C_VERY_LIGHT, line_color=C_ACCENT, line_width=Pt(2))
add_text(slide, Cm(18), Cm(3.5), Cm(14.5), Cm(0.8),
         '일반 도입 대비 기간 단축', font_size=13, bold=True, color=C_ACCENT)

add_shape_text(slide, Cm(18), Cm(4.8), Cm(5), Cm(0.7), '항목', font_size=9, bold=True, color=C_WHITE, fill_color=C_GRAY)
add_shape_text(slide, Cm(23), Cm(4.8), Cm(3.5), Cm(0.7), '일반', font_size=9, bold=True, color=C_WHITE, fill_color=C_GRAY)
add_shape_text(slide, Cm(26.5), Cm(4.8), Cm(3.5), Cm(0.7), 'JB', font_size=9, bold=True, color=C_WHITE, fill_color=C_TEAL)

timeline_data = [
    ("연관관계 분석", "2~4주", "불필요"),
    ("DB 접속 환경", "1~2주", "불필요"),
    ("PII 메타 구성", "2~3주", "3~4일"),
    ("전사 파기+접근이력", "4~8개월", "약 2개월"),
]
for j, (item, normal, jb) in enumerate(timeline_data):
    yy = Cm(5.5 + j * 0.8)
    add_text(slide, Cm(18), yy, Cm(5), Cm(0.8), item, font_size=9, color=C_BLACK, valign='middle')
    add_text(slide, Cm(23), yy, Cm(3.5), Cm(0.8), normal, font_size=9, color=C_GRAY, align=PP_ALIGN.CENTER, valign='middle')
    add_text(slide, Cm(26.5), yy, Cm(3.5), Cm(0.8), jb, font_size=9, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER, valign='middle')

# 구축 타임라인
weeks_data = [
    ("Week 1~3", "X-Scan + X-Purge", ["KeyMap 연동", "X-Scan PII 확정", "Purge Job 구성"], C_RED),
    ("Week 4~6", "X-Audit 구축", ["수집 대상 등록", "DB Audit 연동", "탐지 규칙 설정"], C_ACCENT),
    ("Week 7", "통합 테스트", ["파기/이력 검증", "이상행위 탐지 테스트", "소명 프로세스 검증"], C_ACCENT2),
    ("Week 8", "운영 전환", ["스케줄/대시보드 설정", "담당자 교육", "X-One 거버넌스 완성"], C_GREEN),
]

for i, (week, title, items, color) in enumerate(weeks_data):
    x = Cm(1.5 + i * 8)
    add_rounded_rect(slide, x, Cm(9.5), Cm(7.2), Cm(5.8), fill_color=C_WHITE, line_color=color, line_width=Pt(1.5))
    add_shape_text(slide, x, Cm(9.5), Cm(7.2), Cm(0.8),
                   week, font_size=10, color=C_WHITE, fill_color=color)
    add_text(slide, x, Cm(10.3), Cm(7.2), Cm(0.8),
             title, font_size=12, bold=True, color=color, align=PP_ALIGN.CENTER, valign='middle')
    add_multiline(slide, Emu(x.emu + Cm(0.6).emu), Cm(11.5), Cm(6), Cm(3.5),
                  items, font_size=10, color=C_GRAY)

add_key_point(slide, 'X-Gen 자산 활용으로 도입 기간 60% 이상 단축 — 약 2개월 내 X-One 거버넌스 완성')
add_bottom_bar(slide)
add_page_number(slide, 25)


# ═══════════════════════════════════════════════════════
# ═══ Slide 25: 거버넌스 완성 로드맵 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 24, 'JB우리캐피탈 X-One 거버넌스 완성 로드맵',
                   'X-Gen(완료) → X-Purge + X-Audit(이번 제안) → 전사 거버넌스 체계 완성')

phases = [
    ("완료", "2025", "X-Gen 도입 완료",
     ["연관관계 구축 완료", "KeyMap 확보", "DB 환경 구성 완료", "테스트데이터 자동생성 운영중"], C_GREEN),
    ("이번 제안", "2026 Q2", "X-Scan + X-Purge\n+ X-Audit 추가",
     ["약 6주 구축", "파기 자동화 체계 구축", "접근이력 수집 · 이상행위 탐지", "소명 프로세스 · 통합 대시보드"], C_ACCENT),
    ("목표", "2026 Q3~", "X-One\n거버넌스 완성",
     ["4개 모듈 전체 운영", "통합 CPO 대시보드", "전사 거버넌스 체계 완성", "감사·규제 대응 완벽 증명"], C_TEAL),
]

for i, (label, period, title, items, color) in enumerate(phases):
    x = Cm(1.5 + i * 10.8)
    add_rounded_rect(slide, x, Cm(3.5), Cm(10), Cm(9.5), fill_color=C_WHITE, line_color=color, line_width=Pt(2))

    add_shape_text(slide, Emu(x.emu + Cm(2.5).emu), Cm(3.8), Cm(5), Cm(1),
                   label, font_size=11, bold=True, color=C_WHITE, fill_color=color, shape_type='rounded')

    add_text(slide, x, Cm(5.2), Cm(10), Cm(0.8),
             period, font_size=15, bold=True, color=color, align=PP_ALIGN.CENTER, valign='middle')
    add_text(slide, x, Cm(6.2), Cm(10), Cm(1.8),
             title, font_size=13, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER, valign='middle')
    add_multiline(slide, Emu(x.emu + Cm(1).emu), Cm(8.5), Cm(8), Cm(4),
                  items, font_size=10, color=C_GRAY)

    if i < 2:
        add_text(slide, Cm(11 + i * 10.8), Cm(7), Cm(1.5), Cm(2),
                 '→', font_size=30, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER, valign='middle')

# 모듈 활성화
add_text(slide, Cm(1.5), Cm(13.5), Cm(31), Cm(0.7),
         'X-One 모듈 활성화 현황', font_size=13, bold=True, color=C_BLACK)

module_status = [
    ("X-Scan", "이번", C_ACCENT2),
    ("X-Purge", "이번", C_RED),
    ("X-Gen", "완료 ✓", C_GREEN),
    ("X-Audit", "이번", C_ACCENT),
]
for i, (name, status, color) in enumerate(module_status):
    x = Cm(1.5 + i * 8)
    add_shape_text(slide, x, Cm(14.3), Cm(7.2), Cm(1.5),
                   f'{name}    [{status}]', font_size=12, bold=True, color=C_WHITE, fill_color=color, shape_type='rounded')

add_bottom_bar(slide)
add_page_number(slide, 26)


# ═══════════════════════════════════════════════════════
# ═══ Slide 26: AI 로드맵 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide, C_PURPLE)
add_section_header(slide, 25, 'AI 로드맵 — X-One Intelligence',
                   '규칙 기반 탐지를 넘어, ML/LLM 기반 지능형 개인정보 거버넌스로 진화합니다',
                   header_color=C_PURPLE)

# Current / Phase 1 / Phase 2
phases_ai = [
    ("현재 적용 AI", "", C_GREEN,
     ["LLM 기반 PII 자동 탐지 (운영 중)", "3중 엔진 스코어링 (Meta+Pattern+AI)", "규칙 기반 이상행위 탐지 7가지 (운영 중)"]),
    ("Phase 1", "2026 Q3~Q4", C_ACCENT,
     ["ML 기반 이상행위 탐지 (Isolation Forest)", "사용자 행동 베이스라인 학습", "SQL 위험도 자동 분류 (규칙+LLM)"]),
    ("Phase 2", "2027 Q1~", C_PURPLE,
     ["UEBA (사용자/엔티티 행동 분석)", "자연어 기반 접속기록 분석 (LLM)", "접근권한 최적화 · 컴플라이언스 자동 점검"]),
]

for i, (title, period, color, items) in enumerate(phases_ai):
    x = Cm(1.5 + i * 10.8)
    add_rounded_rect(slide, x, Cm(3.2), Cm(10), Cm(4.2), fill_color=C_WHITE, line_color=color, line_width=Pt(2))
    add_text(slide, Emu(x.emu + Cm(0.5).emu), Cm(3.4), Cm(5), Cm(0.7),
             title, font_size=12, bold=True, color=color, valign='middle')
    if period:
        add_text(slide, Emu(x.emu + Cm(5.5).emu), Cm(3.4), Cm(4), Cm(0.7),
                 period, font_size=10, color=C_GRAY, align=PP_ALIGN.RIGHT, valign='middle')
    add_multiline(slide, Emu(x.emu + Cm(0.5).emu), Cm(4.5), Cm(9), Cm(2.5),
                  items, font_size=10, color=C_BLACK)

    if i < 2:
        add_text(slide, Cm(11 + i * 10.8), Cm(4.5), Cm(1.5), Cm(1.5),
                 '→', font_size=24, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER, valign='middle')

# AI Feature Details
ai_features = [
    ("ML 이상행위 탐지", "국내 최초",
     ["사용자별 접속 패턴 학습 → 베이스라인", "이탈 시 이상 스코어(0~100) 자동 산출",
      "복합적 이상행위 탐지", "벤치마크: Imperva DRA, IBM Guardium"], C_ACCENT),
    ("UEBA", "글로벌 수준",
     ["사용자별 행동 프로파일 카드 생성", "피어그룹 대비 이탈 탐지",
      "리스크 스코어 누적 대시보드", "벤치마크: Securonix, Exabeam"], C_PURPLE),
    ("자연어 로그 분석", "LLM 활용",
     ['"이상 접속 패턴 분석해줘" 자연어 질의', "이상행위 Alert 자연어 설명 자동 생성",
      "월간 분석 보고서 자동 작성", "기존 Privacy-AI LLM 인프라 활용"], C_TEAL),
    ("SQL 위험도 분류", "실시간",
     ["캡처된 SQL 자동 위험도 분류", "Critical: DDL / High: 대량추출",
      "Medium: 민감접근 / Low: 일반업무", "PSM Agent 데이터 즉시 활용"], C_RED),
]

for i, (title, badge, items, color) in enumerate(ai_features):
    x = Cm(1 + i * 8.2)
    add_rounded_rect(slide, x, Cm(8.2), Cm(7.5), Cm(7.5), fill_color=C_WHITE, line_color=color, line_width=Pt(1.5))
    add_shape_text(slide, x, Cm(8.2), Cm(7.5), Cm(1.2),
                   title, font_size=12, bold=True, color=C_WHITE, fill_color=color)
    add_text(slide, x, Cm(9.4), Cm(7.5), Cm(0.7),
             badge, font_size=9, bold=True, color=color, align=PP_ALIGN.CENTER, valign='middle')
    add_multiline(slide, Emu(x.emu + Cm(0.5).emu), Cm(10.5), Cm(6.5), Cm(4.5),
                  items, font_size=9, color=C_GRAY)

add_key_point(slide, '국내 접속기록관리 시장에서 ML/AI 기반 행동 분석을 본격 적용한 솔루션은 아직 없습니다 — X-One이 선도합니다.', color=C_PURPLE)
add_bottom_bar(slide, C_PURPLE)
add_page_number(slide, 27)


# ═══════════════════════════════════════════════════════
# ═══ Slide 27: 레퍼런스 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_top_accent(slide)
add_section_header(slide, 26, '도입 레퍼런스',
                   '금융권 검증을 완료한 엔터프라이즈 솔루션')

refs = [
    ("하나카드", "도입완료", "카드사 개인정보 파기 전사 자동화 체계 구축", C_GREEN),
    ("IM캐피탈", "도입완료", "캐피탈 특화 파기 프로세스, 연관 테이블 완전 파기\n테스트데이터 자동생성 추가 도입 (26년 상반기)", C_GREEN),
    ("수협은행", "26년 도입확정", "은행권 대규모 파기, 이기종 DB 통합 관리", C_ACCENT),
    ("NH농협캐피탈", "26년 도입확정", "개인정보파기 + 테스트데이터 자동생성 통합 구축", C_ACCENT),
]

for i, (name, status, desc, color) in enumerate(refs):
    row = i // 2
    col = i % 2
    x = Cm(1.5 + col * 16)
    y = Cm(3.5 + row * 6.2)

    add_rounded_rect(slide, x, y, Cm(15), Cm(5.2), fill_color=C_WHITE, line_color=color, line_width=Pt(2))

    add_text(slide, Emu(x.emu + Cm(1).emu), Emu(y.emu + Cm(0.5).emu), Cm(8), Cm(1.5),
             name, font_size=20, bold=True, color=C_BLACK, valign='middle')

    add_shape_text(slide, Emu(x.emu + Cm(10).emu), Emu(y.emu + Cm(0.7).emu), Cm(4), Cm(0.9),
                   status, font_size=10, bold=True, color=C_WHITE, fill_color=color, shape_type='rounded')

    add_text(slide, Emu(x.emu + Cm(1).emu), Emu(y.emu + Cm(2.5).emu), Cm(13), Cm(2.5),
             desc, font_size=12, color=C_GRAY, valign='top')

add_bottom_bar(slide)
add_page_number(slide, 28)


# ═══════════════════════════════════════════════════════
# ═══ Slide 28: CTA 마무리 ═══
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, Emu(SLIDE_W), Emu(SLIDE_H), fill_color=C_NAVY)
add_shape(slide, 0, 0, Emu(SLIDE_W), Cm(0.2), fill_color=C_TEAL)
add_shape(slide, Cm(2.5), Cm(2.5), Cm(0.15), Cm(4), fill_color=C_TEAL)

add_text(slide, Cm(3.5), Cm(2.5), Cm(28), Cm(2.5),
         '개인정보 관리,\n이제는 시스템으로 증명하십시오.', font_size=28, bold=True, color=C_WHITE, valign='top',
         line_spacing=1.4)

modules_summary = [
    ("인지    X-Scan", "개인정보가 어디에 있는지 자동으로 발견합니다", C_ACCENT2),
    ("통제    X-Purge", "보유기간 경과 시 완전히 파기하고 증명합니다", C_RED),
    ("감시    X-Audit", "누가 언제 왜 접근했는지 추적하고 소명합니다", C_ACCENT),
    ("활용    X-Gen", "안전한 테스트데이터로 개발 환경을 보호합니다", C_TEAL),
]

for i, (title, desc, color) in enumerate(modules_summary):
    y = Cm(5.8 + i * 1.6)
    add_shape(slide, Cm(3.5), y, Cm(0.15), Cm(1.1), fill_color=color)
    add_text(slide, Cm(4.2), y, Cm(7), Cm(1.1),
             f'✓  {title}', font_size=13, bold=True, color=color, valign='middle')
    add_text(slide, Cm(14), y, Cm(18), Cm(1.1),
             desc, font_size=12, color=RGBColor(0xB0, 0xC4, 0xD8), valign='middle')

add_shape(slide, Cm(3.5), Cm(12.2), Cm(12), Cm(0.05), fill_color=C_TEAL)
add_text(slide, Cm(3.5), Cm(12.6), Cm(15), Cm(1.5),
         'X-One', font_size=40, bold=True, color=C_WHITE, valign='middle')
add_text(slide, Cm(3.5), Cm(14.2), Cm(20), Cm(0.8),
         'All Data.  One Platform.', font_size=15, color=RGBColor(0x8A, 0xB4, 0xD0), valign='middle')
add_text(slide, Cm(3.5), Cm(15.3), Cm(28), Cm(0.8),
         '하나의 플랫폼, 4개 모듈, 전사 개인정보 거버넌스 완성', font_size=13, color=C_WHITE, valign='middle')

# Contact
add_shape(slide, Cm(20), Cm(12.2), Cm(12), Cm(0.05), fill_color=C_TEAL)
add_text(slide, Cm(20), Cm(12.6), Cm(12), Cm(0.8),
         '데이터블록스 (Datablocks)', font_size=15, bold=True, color=C_WHITE, valign='middle')
add_text(slide, Cm(20), Cm(13.6), Cm(12), Cm(0.7),
         '차민석 이사', font_size=13, color=C_WHITE, valign='middle')
add_text(slide, Cm(20), Cm(14.4), Cm(12), Cm(0.7),
         '010-4036-7753', font_size=11, color=RGBColor(0xB0, 0xC4, 0xD8), valign='middle')
add_text(slide, Cm(20), Cm(15.2), Cm(12), Cm(0.7),
         'mscha@datablocks.co.kr', font_size=11, color=RGBColor(0xB0, 0xC4, 0xD8), valign='middle')

# Next Step
add_shape_text(slide, Cm(3.5), Cm(16.8), Cm(28), Cm(1.2),
               'Next Step:  솔루션 데모 시연  →  PoC 범위 확정  →  도입 추진',
               font_size=14, bold=True, color=C_WHITE, fill_color=C_TEAL, shape_type='rounded')

add_page_number(slide, 29)


# ═══ Save ═══
output_path = '/mnt/c/Datablocks/04_Solution/00_공통_제품소개서/710.X-Purge/[X-One 통합 플랫폼] 개인정보 거버넌스 제안서 v8.pptx'
prs.save(output_path)
print(f'Saved: {output_path}')
print(f'Total slides: {len(prs.slides)}')
